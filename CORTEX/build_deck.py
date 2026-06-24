# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ORANGE=RGBColor(0xC4,0x5A,0x1A); BLACK=RGBColor(0x1C,0x19,0x17)
CREAM=RGBColor(0xF5,0xF2,0xEE); WHITE=RGBColor(0xFF,0xFF,0xFF)
GREY=RGBColor(0x6B,0x65,0x60); LGREY=RGBColor(0xB9,0xB2,0xAA)
DARK2=RGBColor(0x2A,0x25,0x22); ORANGE_L=RGBColor(0xE0,0x8A,0x4F)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
blank=prs.slide_layouts[6]
HEAD="Bookman Old Style"; BODY="Calibri"

def slide(bg=WHITE):
    s=prs.slides.add_slide(blank)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s

def box(s,l,t,w,h,lines,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
    first=True
    for ln in lines:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=ln.get("align",align)
        if "space_before" in ln: p.space_before=Pt(ln["space_before"])
        p.space_after=Pt(ln.get("space_after",2))
        if "line" in ln: p.line_spacing=ln["line"]
        runs=ln["runs"] if "runs" in ln else [(ln["t"],ln)]
        for txt,st in runs:
            r=p.add_run(); r.text=txt
            r.font.name=st.get("font",BODY); r.font.size=Pt(st.get("size",16))
            r.font.bold=st.get("bold",False); r.font.italic=st.get("italic",False)
            r.font.color.rgb=st.get("color",BLACK)
    return tb

def rect(s,l,t,w,h,fill,line=None,shadow=False,round=False):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,
        Inches(l),Inches(t),Inches(w),Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line: shp.line.color.rgb=line; shp.line.width=Pt(1)
    else: shp.line.fill.background()
    shp.shadow.inherit=False
    return shp

def circle(s,l,t,d,fill):
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(l),Inches(t),Inches(d),Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb=fill; c.line.fill.background(); c.shadow.inherit=False
    return c

def kicker(s,l,t,text,color=ORANGE):
    box(s,l,t,8,0.3,[{"t":text.upper(),"size":12.5,"bold":True,"color":color,"font":BODY}])

def title(s,l,t,text,w=11.5,color=BLACK,size=33):
    box(s,l,t,w,1.0,[{"t":text,"size":size,"bold":True,"color":color,"font":HEAD,"line":1.0}])

# ---------- helper: simple table as rects ----------
def grid_table(s, l, t, colw, rowh, data, header_fill=ORANGE, head_color=WHITE,
               body_fill=CREAM, alt_fill=WHITE, fs=11, first_bold=True, head_fs=11,
               body_text_color=BLACK):
    y=t
    for ri,row in enumerate(data):
        x=l
        for ci,val in enumerate(row):
            w=colw[ci]
            if ri==0: f=header_fill
            else: f=body_fill if ri%2==1 else alt_fill
            r=rect(s,x,y,w,rowh,f);
            col = head_color if ri==0 else body_text_color
            bold = (ri==0) or (first_bold and ci==0)
            box(s,x+0.08,y,w-0.16,rowh,[{"t":str(val),"size":(head_fs if ri==0 else fs),
                "bold":bold,"color":col,"font":BODY,"align":PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER}],
                anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)
            x+=w
        y+=rowh

# ============ SLIDE 1 — TITLE ============
s=slide(BLACK)
circle(s,-1.5,-1.5,4.5,DARK2)
circle(s,11.2,5.0,4.0,DARK2)
box(s,0.9,2.0,11.5,1.3,[{"t":"CORTEX","size":80,"bold":True,"color":ORANGE,"font":HEAD}])
box(s,0.95,3.45,11.5,0.6,[{"t":"The Business Intelligence Operating System","size":24,"color":WHITE,"font":BODY}])
box(s,0.95,4.15,11.5,0.5,[{"t":"A governed digital brain every organization can install — and grow into.","size":15,"color":LGREY,"font":BODY,"italic":True}])
box(s,0.95,6.35,11.5,0.5,[{"runs":[("Investor Deck  ·  v1.0  ·  June 2026  ·  ",{"size":12,"color":LGREY,"font":BODY}),
    ("Karim Azzaoui · winsym.ai · Kuala Lumpur",{"size":12,"color":ORANGE_L,"font":BODY,"bold":True})]}])

# ============ SLIDE 2 — PROBLEM ============
s=slide(WHITE)
kicker(s,0.9,0.6,"The Problem")
title(s,0.9,0.95,"Every company leaks its own intelligence")
pains=[("Knowledge walks out the door","When a key person leaves, years of context vanish. The business pays to re-learn what it already knew."),
("The same questions, answered 1,000 times","“How do we handle this refund?” “What’s our policy on X?” — answered inconsistently, from memory."),
("Decisions on opinion, not evidence","The loudest voice or latest anecdote wins, because the history is too scattered to consult in the moment."),
("Tools are silos","The CRM doesn’t know the contracts; the contracts don’t know the meeting; the meeting never reaches the work.")]
x=0.9; y=2.25; cw=5.75; ch=2.05; gap=0.35
for i,(h,b) in enumerate(pains):
    cx=x+(i%2)*(cw+gap); cy=y+(i//2)*(ch+gap)
    rect(s,cx,cy,cw,ch,CREAM,round=True)
    circle(s,cx+0.3,cy+0.32,0.5,ORANGE)
    box(s,cx+0.32,cy+0.32,0.46,0.5,[{"t":str(i+1),"size":18,"bold":True,"color":WHITE,"font":HEAD,"align":PP_ALIGN.CENTER}],anchor=MSO_ANCHOR.MIDDLE)
    box(s,cx+1.0,cy+0.28,cw-1.25,0.5,[{"t":h,"size":15.5,"bold":True,"color":BLACK,"font":BODY}])
    box(s,cx+1.0,cy+0.85,cw-1.25,ch-1.0,[{"t":b,"size":12,"color":GREY,"font":BODY,"line":1.08}])
box(s,0.9,6.95,11.5,0.4,[{"t":"Growth without a brain doesn’t scale the company — it scales the chaos.","size":13.5,"italic":True,"bold":True,"color":ORANGE,"font":BODY}])

# ============ SLIDE 3 — SOLUTION ============
s=slide(BLACK)
kicker(s,0.9,0.6,"The Solution",ORANGE_L)
title(s,0.9,0.95,"Not a chatbot. The brain a company runs on.",color=WHITE)
box(s,0.9,1.95,11.5,0.7,[{"t":"CORTEX captures what the business knows, runs how it works, supports its decisions, and acts on a cadence — grounded in stored, sourced knowledge, and it never fabricates.","size":15,"color":LGREY,"font":BODY,"line":1.15}])
pillars=[("Knowledge Graph","Every decision, customer, contract and risk is a node; relationships are typed edges. The brain reasons, not just retrieves."),
("Governance / RBAC","Sensitivity on the data, clearance on the user, enforced server-side. The CEO and the front desk see different views of one system."),
("Grounded Trust Layer","Answers only from what’s stored — with sources, a confidence band, and a named human approver for high-stakes output.")]
x=0.9; cw=3.74; gap=0.16; y=3.0
for i,(h,b) in enumerate(pillars):
    cx=x+i*(cw+gap)
    rect(s,cx,y,cw,3.0,DARK2,round=True)
    circle(s,cx+0.35,y+0.35,0.55,ORANGE)
    box(s,cx+0.35,y+0.35,0.55,0.55,[{"t":str(i+1),"size":20,"bold":True,"color":WHITE,"font":HEAD,"align":PP_ALIGN.CENTER}],anchor=MSO_ANCHOR.MIDDLE)
    box(s,cx+0.35,y+1.1,cw-0.7,0.6,[{"t":h,"size":16,"bold":True,"color":ORANGE_L,"font":BODY}])
    box(s,cx+0.35,y+1.75,cw-0.7,1.1,[{"t":b,"size":12,"color":LGREY,"font":BODY,"line":1.12}])
box(s,0.9,6.5,11.5,0.5,[{"t":"3 pillars · 8 functional layers · 1 continuous improvement loop (the CORTEX Method)","size":13,"italic":True,"color":WHITE,"font":BODY,"align":PP_ALIGN.CENTER}])

# ============ SLIDE 4 — CATEGORY / WHY NOW ============
s=slide(WHITE)
kicker(s,0.9,0.6,"The Category")
title(s,0.9,0.95,"We’re creating a category: BIOS")
box(s,0.9,1.9,7.0,1.2,[{"t":"Business Intelligence Operating System — the governed layer that sits above every tool, person and decision in the company. Search, chat, CRM and agents are features of it, not alternatives to it.","size":15,"color":GREY,"font":BODY,"line":1.18}])
# why now stats
stats=[("~80%","of enterprises deploying generative AI by 2026"),
("$12B → $80B","SEA AI market, 2026 → 2031"),
("~68%","Malaysia professional AI adoption"),
("14–18%","KM software market CAGR")]
x=0.9; y=3.5; cw=2.85; gap=0.2
for i,(n,l) in enumerate(stats):
    cx=x+i*(cw+gap)
    rect(s,cx,y,cw,2.4,CREAM,round=True)
    box(s,cx+0.15,y+0.45,cw-0.3,0.9,[{"t":n,"size":30,"bold":True,"color":ORANGE,"font":HEAD,"align":PP_ALIGN.CENTER}])
    box(s,cx+0.2,y+1.5,cw-0.4,0.8,[{"t":l,"size":12,"color":GREY,"font":BODY,"align":PP_ALIGN.CENTER,"line":1.08}])
box(s,0.9,6.35,11.5,0.5,[{"t":"High intent, low maturity, weak governance — the perfect moment for a governed brain, starting in Southeast Asia.","size":13.5,"italic":True,"bold":True,"color":BLACK,"font":BODY}])

# ============ SLIDE 5 — HOW IT WORKS / 8 LAYERS ============
s=slide(WHITE)
kicker(s,0.9,0.6,"The Product")
title(s,0.9,0.95,"One brain, eight functions")
layers=["Knowledge","Operations","Marketing & Content","Sales","Customer","Decision","Intelligence & Risk","Automation"]
subs=["single source of truth","runs the work","on-brand at scale","every rep is your best rep","one memory per relationship","score high-stakes choices","live radar + crisis playbook","runs on a cadence"]
x=0.9; y=2.1; cw=2.85; ch=1.55; gx=0.2; gy=0.25
for i,(h,b) in enumerate(zip(layers,subs)):
    cx=x+(i%4)*(cw+gx); cy=y+(i//4)*(ch+gy)
    rect(s,cx,cy,cw,ch,CREAM if (i//4+i%4)%2==0 else WHITE,line=LGREY,round=True)
    circle(s,cx+0.22,cy+0.25,0.42,ORANGE)
    box(s,cx+0.22,cy+0.25,0.42,0.42,[{"t":str(i+1),"size":14,"bold":True,"color":WHITE,"font":HEAD,"align":PP_ALIGN.CENTER}],anchor=MSO_ANCHOR.MIDDLE)
    box(s,cx+0.78,cy+0.22,cw-0.9,0.5,[{"t":h,"size":13.5,"bold":True,"color":BLACK,"font":BODY,"line":0.95}])
    box(s,cx+0.2,cy+0.92,cw-0.4,0.55,[{"t":b,"size":11,"color":GREY,"font":BODY,"italic":True,"line":1.0}])
box(s,0.9,6.35,11.5,0.5,[{"t":"A small business switches on 3 layers; an enterprise runs all 8. Same platform, different views.","size":13,"italic":True,"color":ORANGE,"font":BODY,"bold":True}])

# ============ SLIDE 6 — MATURITY LADDER ============
s=slide(BLACK)
kicker(s,0.9,0.6,"The Engine of Pricing & Sales",ORANGE_L)
title(s,0.9,0.95,"Clients buy their way up a 5-level ladder",color=WHITE)
levels=[("L1","Documentation","Is it written down?"),
("L2","Knowledge Mgmt","Can anyone find & trust the answer?"),
("L3","Process Intelligence","Can the brain run our work?"),
("L4","Decision Intelligence","Does it make us decide better?"),
("L5","Operating System","Does the business run on it?")]
x=0.9; y=2.4; cw=2.3; gap=0.13
for i,(lv,nm,q) in enumerate(levels):
    cx=x+i*(cw+gap); h=2.0+i*0.45
    cy=y+(2.3-i*0.45)
    shade=RGBColor(0x3a-i*4,0x2a,0x20+i*6) if False else DARK2
    rect(s,cx,cy,cw,h,DARK2,round=True)
    box(s,cx,cy+0.2,cw,0.6,[{"t":lv,"size":24,"bold":True,"color":ORANGE,"font":HEAD,"align":PP_ALIGN.CENTER}])
    box(s,cx+0.12,cy+0.85,cw-0.24,0.6,[{"t":nm,"size":13,"bold":True,"color":WHITE,"font":BODY,"align":PP_ALIGN.CENTER,"line":0.95}])
    box(s,cx+0.12,cy+1.45,cw-0.24,h-1.5,[{"t":q,"size":10.5,"color":LGREY,"font":BODY,"italic":True,"align":PP_ALIGN.CENTER,"line":1.0}])
box(s,0.9,6.95,11.5,0.4,[{"t":"The iron rule: never automate (L5) on top of unverified knowledge (L1). The maturity gap is the value case in every sale.","size":12,"italic":True,"color":WHITE,"font":BODY,"align":PP_ALIGN.CENTER}])

# ============ SLIDE 7 — BUSINESS MODEL ============
s=slide(WHITE)
kicker(s,0.9,0.6,"Business Model")
title(s,0.9,0.95,"Build once. Configure per industry. Sell many.")
cols=[("Platform","The 8-layer brain, sold as recurring SaaS across 4 tiers.","Recurring revenue"),
("Implementation","The CORTEX Method captures the client’s knowledge and stands up their brain.","High-margin wedge"),
("Industry Templates","Reusable per-vertical packs — vocabulary, scoring, compliance, starter SOPs.","The scalability engine")]
x=0.9; y=2.2; cw=3.74; gap=0.16
for i,(h,b,tag) in enumerate(cols):
    cx=x+i*(cw+gap)
    rect(s,cx,y,cw,3.2,CREAM,round=True)
    box(s,cx+0.3,y+0.3,cw-0.6,0.6,[{"t":h,"size":18,"bold":True,"color":ORANGE,"font":HEAD}])
    box(s,cx+0.3,y+1.0,cw-0.6,1.6,[{"t":b,"size":13,"color":GREY,"font":BODY,"line":1.15}])
    box(s,cx+0.3,y+2.55,cw-0.6,0.5,[{"t":tag.upper(),"size":11.5,"bold":True,"color":BLACK,"font":BODY}])
box(s,0.9,5.9,11.5,0.9,[{"t":"Services-led to learn the verticals and fund the build → platform-led recurring revenue as templates mature. A brain must be trusted before it’s bought; new verticals are new templates, not new engineering.","size":13,"italic":True,"color":BLACK,"font":BODY,"line":1.15}])

# ============ SLIDE 8 — OFFER LADDER ============
s=slide(WHITE)
kicker(s,0.9,0.6,"The Offer")
title(s,0.9,0.95,"Four tiers, one upgrade path")
grid_table(s,0.9,2.1,[1.65,2.15,2.15,2.15],0.52,
 [["","STARTER","GROWTH","PROFESS."],
  ["Target","Solo / micro","SME (10–250)","Mid-market"],
  ["Maturity","L1 → 2","L2 → 3","L3 → 4"],
  ["Layers","2–3","4–5","6–7"],
  ["Automations","Read-only","Schedules","Workflows"],
  ["Deployment","Multi-tenant","Multi-tenant","MT/isolated"]],fs=11)
# Enterprise highlighted column
rect(s,9.35,2.1,3.05,3.12,BLACK,round=True)
ent=[("ENTERPRISE",16,WHITE,True),("2.5K–10K+ employees",11.5,LGREY,False),("Level 4 → 5",11.5,ORANGE_L,True),
("All 8 lobes",11.5,WHITE,False),("Full cadence + triggers",11.5,WHITE,False),("Single-tenant / VPC / on-prem",11.5,WHITE,False)]
yy=2.28
for t,sz,cl,bd in ent:
    box(s,9.6,yy,2.65,0.5,[{"t":t,"size":sz,"bold":bd,"color":cl,"font":(HEAD if t=="ENTERPRISE" else BODY)}]); yy+=0.5
box(s,0.9,5.55,11.5,1.2,[{"runs":[("Add-ons monetize the variable layer: ",{"size":13,"bold":True,"color":BLACK,"font":BODY}),
("extra industry templates, additional brand-brains (agencies), premium integrations, compliance packs.",{"size":13,"color":GREY,"font":BODY})],"line":1.15}])

# ============ SLIDE 9 — PRICING ============
s=slide(WHITE)
kicker(s,0.9,0.6,"Pricing")
title(s,0.9,0.95,"Priced on value & maturity — three regions")
box(s,0.9,1.85,11.5,0.4,[{"t":"Year-1 total (setup + 12 months). SEA is the anchor; EU and Premium are variants.","size":12.5,"color":GREY,"font":BODY,"italic":True}])
grid_table(s,0.9,2.4,[2.4,2.35,2.35,2.3,2.1],0.55,
 [["Region","STARTER","GROWTH","PROFESSIONAL","ENTERPRISE"],
  ["Malaysia / SEA","RM 20.4K","RM 60K","RM 141K","RM 336K+"],
  ["Europe / DACH","€ 7.9K","€ 24.8K","€ 64K","€ 167K+"],
  ["Premium (USD)","$ 10.2K","$ 36K","$ 107K","$ 270K+"]],fs=13,head_fs=12)
box(s,0.9,5.0,11.5,1.6,[{"runs":[("The wedge: ",{"size":14,"bold":True,"color":ORANGE,"font":BODY}),
("our Growth tier (~USD 13K Year-1) lands below the ~$60K–108K floors of Glean and ChatGPT Enterprise — while delivering far more than search or chat. Enterprise still undercuts a Palantir engagement by an order of magnitude. We win on value-per-dollar at the mid-market, not per-seat price wars.",{"size":14,"color":BLACK,"font":BODY})],"line":1.2}])

# ============ SLIDE 10 — COMPETITION ============
s=slide(BLACK)
kicker(s,0.9,0.6,"Competition",ORANGE_L)
title(s,0.9,0.95,"Everyone owns a slice. Nobody owns the brain.",color=WHITE)
grid_table(s,0.9,2.15,[3.1,4.5,3.9],0.5,
 [["Player","Indicative 2026 price","Gap vs. CORTEX"],
  ["Palantir Foundry/AIP","Custom; 6–7 figures","Unreachable for SME/mid-market"],
  ["Glean","~$50/user; ~$60K+ ACV floor","Answers only; no execution/memory"],
  ["ChatGPT Enterprise","~$50–60/user; ~$108K/yr floor","Generic; can fabricate; no governed graph"],
  ["MS Copilot / Notion AI","$30 / $20 per user","No business graph, weak governance"],
  ["Salesforce Agentforce","$2/conv · ~$125/user","Locked-in; agents on ungoverned base"]],
 header_fill=ORANGE,body_fill=DARK2,alt_fill=RGBColor(0x33,0x2d,0x29),head_color=WHITE,fs=11.5,
 body_text_color=CREAM)
box(s,0.9,5.85,11.5,1.1,[{"runs":[("Our moat: ",{"size":14,"bold":True,"color":ORANGE_L,"font":BODY}),
("the governance + trust layer, the compounding knowledge graph + memory, the proprietary template library, and the repeatable CORTEX Method. Generic AI can’t copy the trust layer without rebuilding its core.",{"size":14,"color":WHITE,"font":BODY})],"line":1.18}])

# ============ SLIDE 11 — TRACTION / LINEAGE ============
s=slide(WHITE)
kicker(s,0.9,0.6,"Why This Is Credible")
title(s,0.9,0.95,"Productized from a system that already works")
box(s,0.9,2.0,6.0,3.4,[
 {"t":"CORTEX is the domain-neutral extraction of the MIPP PIOS — a real, governed “leadership brain” built for a national political organization.","size":15,"color":BLACK,"font":BODY,"line":1.2,"space_after":10},
 {"t":"It solved the hardest version of a universal problem: capture what the people at the top know, prevent the AI from making things up, and control who sees what — in a domain where one fabricated fact causes real damage.","size":14,"color":GREY,"font":BODY,"line":1.2,"space_after":10},
 {"t":"Remove the politics and what remains — governance, grounded retrieval, memory, the knowledge graph, the maturity model, the verification layer — is already industry-neutral.","size":14,"color":GREY,"font":BODY,"line":1.2}])
rect(s,7.4,2.0,5.0,3.9,CREAM,round=True)
box(s,7.7,2.25,4.4,0.5,[{"t":"Already in hand","size":15,"bold":True,"color":ORANGE,"font":BODY}])
for i,t in enumerate(["Proven governed architecture (graph + RBAC + grounded RAG)",
"The CORTEX Method delivery system","A framework ready to template across verticals",
"Founder with Google + Meta Partner credibility, KL-based"]):
    circle(s,7.7,2.95+i*0.72,0.34,ORANGE)
    box(s,7.75,2.95+i*0.72,0.34,0.34,[{"t":"✓","size":13,"bold":True,"color":WHITE,"font":BODY,"align":PP_ALIGN.CENTER}],anchor=MSO_ANCHOR.MIDDLE)
    box(s,8.2,2.93+i*0.72,3.95,0.7,[{"t":t,"size":12.5,"color":BLACK,"font":BODY,"line":1.05}],anchor=MSO_ANCHOR.MIDDLE)
box(s,0.9,6.4,11.5,0.5,[{"t":"The political build stays a confidential, separate implementation. CORTEX is the company.","size":12,"italic":True,"color":GREY,"font":BODY}])

# ============ SLIDE 12 — GTM ============
s=slide(WHITE)
kicker(s,0.9,0.6,"Go-to-Market")
title(s,0.9,0.95,"The diagnostic is the pitch")
steps=[("Diagnose","Place the prospect on the maturity ladder — the gap is the value case."),
("Demo","A live, grounded, fully-sourced answer — then a deliberate “I don’t know.” That closes regulated buyers."),
("Land","Start with the 1–2 layers that hurt most; prove ROI fast at Level 2."),
("Expand","Sell the climb: more layers, departments, integrations. NRR ~115%.")]
x=0.9; y=2.3; cw=2.85; gap=0.2
for i,(h,b) in enumerate(steps):
    cx=x+i*(cw+gap)
    circle(s,cx,y,0.7,ORANGE)
    box(s,cx,y,0.7,0.7,[{"t":str(i+1),"size":26,"bold":True,"color":WHITE,"font":HEAD,"align":PP_ALIGN.CENTER}],anchor=MSO_ANCHOR.MIDDLE)
    box(s,cx,y+0.95,cw-0.1,0.5,[{"t":h,"size":17,"bold":True,"color":BLACK,"font":BODY}])
    box(s,cx,y+1.5,cw-0.1,2.0,[{"t":b,"size":12.5,"color":GREY,"font":BODY,"line":1.15}])
box(s,0.9,5.9,11.5,1.0,[{"runs":[("Beachhead: ",{"size":13.5,"bold":True,"color":ORANGE,"font":BODY}),
("regulated verticals (law firms, clinics) where the trust layer sells itself — KL/Malaysia first, then SEA, then DACH via the German-language advantage. Win 2–3 references per vertical, harden the template, repeat.",{"size":13.5,"color":BLACK,"font":BODY})],"line":1.18}])

# ============ SLIDE 13 — ROADMAP ============
s=slide(BLACK)
kicker(s,0.9,0.6,"Roadmap",ORANGE_L)
title(s,0.9,0.95,"Framework → Services → Platform → SaaS",color=WHITE)
phases=[("Phase 1","Done-for-you Brain Setup","Services-led pilots, one vertical. Revenue + references."),
("Phase 2","Brain as a Service","Productize: graph + RBAC + grounded RAG + trust layer."),
("Phase 3","Managed Operating System","Switch on automation + integrations for earned maturity."),
("Phase 4","Own SaaS Platform","Self-serve floor + enterprise ceiling + template marketplace.")]
x=0.9; y=2.4; cw=2.85; gap=0.2
for i,(p,h,b) in enumerate(phases):
    cx=x+i*(cw+gap)
    rect(s,cx,y,cw,3.0,DARK2,round=True)
    box(s,cx+0.25,y+0.25,cw-0.5,0.5,[{"t":p,"size":15,"bold":True,"color":ORANGE,"font":HEAD}])
    box(s,cx+0.25,y+0.85,cw-0.5,0.9,[{"t":h,"size":14.5,"bold":True,"color":WHITE,"font":BODY,"line":1.0}])
    box(s,cx+0.25,y+1.85,cw-0.5,1.0,[{"t":b,"size":11.5,"color":LGREY,"font":BODY,"line":1.12}])
box(s,0.9,6.0,11.5,0.6,[{"t":"Sequencing is the strategy: trust, then revenue, then product, then scale.","size":13.5,"italic":True,"color":WHITE,"font":BODY,"align":PP_ALIGN.CENTER}])

# ============ SLIDE 14 — FINANCIALS ============
s=slide(WHITE)
kicker(s,0.9,0.6,"Financials — Base Case")
title(s,0.9,0.95,"A capital-efficient path to $2.6M revenue")
grid_table(s,0.9,2.15,[3.5,2.7,2.7,2.6],0.52,
 [["USD","Year 1","Year 2","Year 3"],
  ["Total revenue","$214K","$921K","$2.59M"],
  ["Recurring (ARR)","$148K","$670K","$1.95M"],
  ["Gross margin","72%","73%","74%"],
  ["EBITDA","($25K)","$154K","$765K"],
  ["Active clients","16","76","243"]],fs=13,head_fs=12)
# stat callouts right
calls=[("$1.95M","Year-3 ARR"),("~74%","Gross margin"),("~115%","Net rev. retention")]
# place below
yy=5.5
for i,(n,l) in enumerate(calls):
    cx=0.9+i*3.85
    rect(s,cx,yy,3.6,1.3,CREAM,round=True)
    box(s,cx+0.2,yy+0.18,3.2,0.7,[{"t":n,"size":28,"bold":True,"color":ORANGE,"font":HEAD,"align":PP_ALIGN.CENTER}])
    box(s,cx+0.2,yy+0.92,3.2,0.4,[{"t":l,"size":12,"color":GREY,"font":BODY,"align":PP_ALIGN.CENTER}])

# ============ SLIDE 15 — THE ASK ============
s=slide(BLACK)
circle(s,10.8,-1.4,4.0,DARK2)
kicker(s,0.9,0.7,"The Ask",ORANGE_L)
title(s,0.9,1.1,"Raising to fund the platform",color=WHITE,size=34)
box(s,0.9,2.2,7.2,2.6,[
 {"t":"Bootstrap Phase 1 from winsym.ai cash flow + a USD 150–300K angel/grant round.","size":16,"color":WHITE,"font":BODY,"line":1.2,"space_after":10},
 {"t":"Use of funds: first 2 engineers, a sales hire, and the platform MVP (Phase 2).","size":14.5,"color":LGREY,"font":BODY,"line":1.2,"space_after":10},
 {"t":"Convert to a priced seed (~$1.5–3M) once 2–3 reference clients and one hardened vertical template are in place.","size":14.5,"color":LGREY,"font":BODY,"line":1.2}])
rect(s,8.4,2.2,4.0,3.5,DARK2,round=True)
box(s,8.7,2.45,3.5,0.5,[{"t":"Why invest now","size":15,"bold":True,"color":ORANGE_L,"font":BODY}])
for i,t in enumerate(["New category (BIOS) with a defensible governance moat","Services-to-SaaS margin story (40% → 80%+)","SEA-to-global expansion path","Proven architecture, not a science project"]):
    circle(s,8.7,3.05+i*0.62,0.3,ORANGE)
    box(s,9.15,3.0+i*0.62,3.0,0.62,[{"t":t,"size":12,"color":WHITE,"font":BODY,"line":1.0}],anchor=MSO_ANCHOR.MIDDLE)
box(s,0.9,6.4,11.5,0.5,[{"t":"work@edgegraphics.co  ·  winsym.ai  ·  +60 11 2448 7055","size":13,"color":ORANGE_L,"font":BODY,"bold":True}])

# ============ SLIDE 16 — CLOSE ============
s=slide(BLACK)
circle(s,-1.6,4.5,4.5,DARK2)
box(s,0.9,2.6,11.5,1.0,[{"t":"CORTEX","size":60,"bold":True,"color":ORANGE,"font":HEAD}])
box(s,0.95,3.85,11.5,0.8,[{"t":"Give every organization a brain — governed, trustworthy, and built to grow with them.","size":18,"color":WHITE,"font":BODY,"italic":True,"line":1.15}])
box(s,0.95,5.2,11.5,0.4,[{"t":"Winning systems, backed by a brain.","size":14,"color":ORANGE_L,"font":BODY,"bold":True}])

prs.save("/sessions/brave-zen-fermat/mnt/Winsym Coaching/CORTEX/CORTEX_Investor_Deck.pptx")
print("deck saved", len(prs.slides.__iter__.__self__._sldIdLst))
